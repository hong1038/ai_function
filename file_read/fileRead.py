"""
fileRead.py

이 파일은 다양한 확장자의 파일에서 텍스트를 추출하는 함수들을 제공합니다.
각 함수는 특정 파일 형식을 처리하며, 독립적으로 호출될 수 있습니다.

지원되는 파일 확장자:
- .txt
- .csv
- .hwp
- .pdf
- .docx
- .xls, .xlsx
- .pptx
- .rtf
- .json
- .xml
- .epub
- .html, .htm
- .log
- .yaml, .yml

사용 방법:
1. 해당 파일을 import하여 특정 파일 형식에 맞는 함수를 호출하십시오.
2. 각 함수는 `file_path` 매개변수를 입력으로 받아 텍스트 데이터를 반환합니다.
"""

import csv
from olefile import OleFileIO
from PyPDF2 import PdfReader
import docx
import openpyxl
from pptx import Presentation
import striprtf
import json
from xml.etree import ElementTree as ET
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup
import yaml


def read_txt_file(file_path):
    """텍스트(.txt) 파일에서 텍스트를 읽어옵니다."""
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    return content


def read_csv_file(file_path):
    """CSV(.csv) 파일에서 텍스트를 읽어옵니다."""
    with open(file_path, 'r', encoding='utf-8') as file:
        reader = csv.reader(file)
        content = [" ".join(row) for row in reader]
    return " ".join(content)


def read_hwp_file(file_path):
    """한글(.hwp) 파일에서 텍스트를 읽어옵니다."""
    with OleFileIO(file_path) as ole:
        content = ole.openstream("BodyText").read().decode("utf-16le", errors="ignore")
    return content


def read_pdf_file(file_path):
    """PDF(.pdf) 파일에서 텍스트를 읽어옵니다."""
    reader = PdfReader(file_path)
    content = " ".join([page.extract_text() for page in reader.pages])
    return content


def read_docx_file(file_path):
    """Word(.docx) 파일에서 텍스트를 읽어옵니다."""
    doc = docx.Document(file_path)
    content = " ".join([paragraph.text for paragraph in doc.paragraphs])
    return content


def read_excel_file(file_path):
    """Excel(.xls, .xlsx) 파일에서 텍스트를 읽어옵니다."""
    workbook = openpyxl.load_workbook(file_path, data_only=True)
    content = []
    for sheet in workbook:
        for row in sheet.iter_rows(values_only=True):
            row_content = " ".join(map(str, row))
            content.append(row_content)
    return " ".join(content)


def read_pptx_file(file_path):
    """PowerPoint(.pptx) 파일에서 텍스트를 읽어옵니다."""
    presentation = Presentation(file_path)
    content = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                content.append(shape.text)
    return " ".join(content)


def read_rtf_file(file_path):
    """RTF(.rtf) 파일에서 텍스트를 읽어옵니다."""
    with open(file_path, 'r', encoding='utf-8') as file:
        content = striprtf.rtf_to_text(file.read())
    return content


def read_json_file(file_path):
    """JSON(.json) 파일에서 텍스트를 읽어옵니다."""
    with open(file_path, 'r', encoding='utf-8') as file:
        data = json.load(file)
    return json.dumps(data, indent=2)


def read_xml_file(file_path):
    """XML(.xml) 파일에서 텍스트를 읽어옵니다."""
    tree = ET.parse(file_path)
    root = tree.getroot()
    content = ET.tostring(root, encoding='utf-8').decode('utf-8')
    return content


def read_epub_file(file_path):
    """EPUB(.epub) 파일에서 텍스트를 읽어옵니다."""
    book = epub.read_epub(file_path)
    content = []
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            soup = BeautifulSoup(item.get_content(), 'html.parser')
            content.append(soup.get_text())
    return " ".join(content)


def read_html_file(file_path):
    """HTML(.html, .htm) 파일에서 텍스트를 읽어옵니다."""
    with open(file_path, 'r', encoding='utf-8') as file:
        soup = BeautifulSoup(file, 'html.parser')
    return soup.get_text()


def read_log_file(file_path):
    """로그(.log) 파일에서 텍스트를 읽어옵니다."""
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    return content


def read_yaml_file(file_path):
    """YAML(.yaml, .yml) 파일에서 텍스트를 읽어옵니다."""
    with open(file_path, 'r', encoding='utf-8') as file:
        data = yaml.safe_load(file)
    return str(data)
